from flask import Flask, request, jsonify
import os
import csv
from PyPDF2 import PdfReader
import openpyxl
from pydub import AudioSegment
from docx import Document
from pptx import Presentation
import json
import xml.etree.ElementTree as ET
from youtube_transcript_api import YouTubeTranscriptApi
from bs4 import BeautifulSoup
import requests
import threading
import time
from moviepy.editor import VideoFileClip


app = Flask(__name__)
API_KEY = "placeholderAPIKey"
AUDIO_EXTENSIONS = ["flac", "mp3", "mp4", "mpeg", "mpga", "m4a", "ogg", "wav", "webm"]


@app.route('/convert', methods=['POST'])
def convert():
    provided_api_key = request.headers.get('x-api-key')
    if not provided_api_key or provided_api_key != API_KEY:
        return jsonify(error="Invalid or missing API key"), 401
    
    if 'file' not in request.files and 'url' not in request.form:
        return jsonify(error="No file or URL provided"), 400
    
    if 'webhookURL' not in request.form:
        return jsonify(error="No webhookURL provided"), 400

    if 'uniqueID' not in request.form:
        return jsonify(error="No uniqueID provided"), 400
    
    webhook_url = request.form['webhookURL']
    unique_id = request.form['uniqueID']

    # Check for OpenAI API key if the file type is an audio type
    if 'file' in request.files:
        file = request.files['file']
        _, extension = os.path.splitext(file.filename)
        extension = extension[1:].lower()  # Removes the dot and make lowercase
        if extension in AUDIO_EXTENSIONS:
            openai_api_key = request.headers.get('OpenAIAPIKey')
            if not openai_api_key:
                return jsonify(error="OpenAIAPIKey is required for audio files"), 400
        else:
            openai_api_key = None  # Not an audio file, so no key needed

    # Asynchronous handling
    if 'file' in request.files:
        file = request.files['file']
        if file.filename == '':
            return jsonify(error="No selected file"), 400

        filename = os.path.join("/tmp", file.filename)
        file.save(filename)
        thread = threading.Thread(target=process_file, args=(filename, webhook_url, unique_id, openai_api_key))
        thread.start()
        return jsonify(status="Received file. Processing..."), 202

    if 'url' in request.form:
        url = request.form['url']
        thread = threading.Thread(target=process_url, args=(url, webhook_url, unique_id))
        thread.start()
        return jsonify(status="Received URL. Processing..."), 202

# Processing the file asynchronously
def process_file(filename, webhook_url, unique_id, openai_api_key=None):
    start_time = time.time()
    text, error = extract_text_from_file(filename, openai_api_key)
    end_time = time.time()
    processing_time = round(end_time - start_time, 2)  # in seconds
    _, extension = os.path.splitext(filename)
    filetype = extension[1:]  # Removes the dot
    os.remove(filename)
    if error:
        payload = {"error": error, "uniqueID": unique_id, "processingTime": processing_time, "filetype": filetype}
    else:
        payload = {"text": text, "uniqueID": unique_id, "processingTime": processing_time, "filetype": filetype}
    send_webhook(webhook_url, payload)

def process_url(url, webhook_url, unique_id):
    start_time = time.time()
    text, error = extract_text_from_url(url)
    end_time = time.time()
    processing_time = round(end_time - start_time, 2)  # in seconds
    
    # Determine the filetype for the URL
    if "youtube.com/watch?v=" in url or "youtu.be/" in url:
        filetype = "youtube"
    else:
        filetype = "url"
    
    if error:
        payload = {"error": error, "uniqueID": unique_id, "processingTime": processing_time, "filetype": filetype}
    else:
        payload = {"text": text, "uniqueID": unique_id, "processingTime": processing_time, "filetype": filetype}
    send_webhook(webhook_url, payload)


# Sending the result to the webhook
def send_webhook(url, payload):
    try:
        response = requests.post(url, json=payload)
        response.raise_for_status()
    except Exception as e:
        print(f"Failed to send webhook. Error: {e}")

def extract_text_from_file(filename, openai_api_key):
    _, extension = os.path.splitext(filename)
    extension = extension.lower()  # Convert extension to lowercase for consistent comparison
    try:
        if extension[1:] in AUDIO_EXTENSIONS:
            return extract_from_audio(filename, openai_api_key), None
        elif extension == '.pdf':
            return extract_from_pdf(filename), None
        elif extension == '.docx':
            return extract_from_word(filename), None
        elif extension == '.xlsx':
            return extract_from_excel(filename), None
        elif extension == '.csv':
            return extract_from_csv(filename), None
        elif extension == '.txt':
            return extract_from_txt(filename), None
        elif extension == '.pptx':
            return extract_from_pptx(filename), None
        elif extension == '.html':
            return extract_from_html(filename), None
        elif extension == '.xml':
            return extract_from_xml(filename), None
        elif extension == '.json':
            return extract_from_json(filename), None
        else:
            return None, "Unsupported filetype"
    except Exception as e:
        return None, str(e)
    
def split_audio(filename):
    # Load the audio file
    audio = AudioSegment.from_file(filename, format="mp3")  # You might need to adjust the format based on your actual file
    length_audio = len(audio)
    chunks = []

    # 60 seconds * 1000ms/sec = 60000 ms
    one_minute = 60 * 1000

    # Split audio into 60-second chunks
    for i in range(0, length_audio, one_minute):
        chunk = audio[i:i + one_minute]
        chunks.append(chunk)

    return chunks

def extract_from_audio(filename, openai_api_key):
    # Define the URL for the OpenAI API within the function's scope
    url = "https://api.openai.com/v1/audio/transcriptions"
    
    headers = {
        "Authorization": f"Bearer {openai_api_key}"
    }

    # Convert video to audio if the file is a video
    _, extension = os.path.splitext(filename)
    if extension[1:].lower() in ["mp4", "mpeg", "webm"]:
        video = VideoFileClip(filename)
        audio_filename = f"/tmp/temp_audio.mp3"
        video.audio.write_audiofile(audio_filename, codec='mp3')
        filename = audio_filename  # Update filename to point to the extracted audio

    # Split the audio into 60-second chunks
    audio_chunks = split_audio(filename)

    transcribed_texts = []

    for index, chunk in enumerate(audio_chunks):
        # Save each chunk to a unique temporary file
        chunk_filename = f"/tmp/temp_chunk_{index}.mp3"  # Adjust format if needed
        chunk.export(chunk_filename, format="mp3")

        # Check chunk size
        chunk_size = os.path.getsize(chunk_filename)
        max_size = 25 * 1024 * 1024  # 25MB in bytes
        if chunk_size > max_size:
            # Clean up any temporary files and then raise the error
            os.remove(chunk_filename)
            if 'audio_filename' in locals():
                os.remove(audio_filename)  # Clean up the extracted audio file if it exists
            raise Exception(f"Chunk {index} exceeds the maximum allowed size of 25MB.")

        # Transcribe the chunk
        with open(chunk_filename, 'rb') as audio_file:
            files = {"file": (os.path.basename(chunk_filename), audio_file)}
            data = {"model": "whisper-1"}

            response = requests.post(url, headers=headers, files=files, data=data)
            response_data = response.json()

            # Clean up the temporary file after processing
            os.remove(chunk_filename)

            if "text" in response_data:
                transcribed_texts.append(response_data["text"])
            else:
                # Providing a more detailed error message
                error_message = response_data.get("error", {}).get("message", "Unknown error")
                raise Exception(f"Failed to transcribe chunk {index}. OpenAI error: {error_message}")

    return " ".join(transcribed_texts)

def extract_from_pdf(filename):
    with open(filename, 'rb') as pdf_file:
        reader = PdfReader(pdf_file)
        text = " ".join(page.extract_text() for page in reader.pages)
    return text

def extract_from_word(filename):
    doc = Document(filename)
    return " ".join(paragraph.text for paragraph in doc.paragraphs)

def extract_from_excel(filename):
    wb = openpyxl.load_workbook(filename)
    all_text = ""
    for sheet in wb:
        for row in sheet.iter_rows():
            for cell in row:
                all_text += str(cell.value) + " "
    return all_text

def extract_from_csv(filename):
    all_text = ""
    with open(filename, 'r') as csv_file:
        reader = csv.reader(csv_file)
        for row in reader:
            all_text += " ".join(row) + " "
    return all_text


def extract_from_txt(filename):
    with open(filename, 'r') as txt_file:
        return txt_file.read()

def extract_from_pptx(filename):
    prs = Presentation(filename)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def extract_from_html(filename):
    with open(filename, 'r', encoding='utf-8') as html_file:
        soup = BeautifulSoup(html_file, 'html.parser')
        return " ".join(soup.stripped_strings)

def extract_from_xml(filename):
    tree = ET.parse(filename)
    root = tree.getroot()
    texts = [elem.text for elem in root.iter() if elem.text]
    return " ".join(texts)

def extract_from_json(filename):
    with open(filename, 'r', encoding='utf-8') as json_file:
        data = json.load(json_file)
        return json.dumps(data)

def extract_text_from_url(url):
    try:
        # Check if the URL is a standard YouTube video URL
        if "youtube.com/watch?v=" in url:
            video_id = url.split("v=")[1].split("&")[0]
            return extract_youtube_transcript(video_id)
        # Check if the URL is a shortened YouTube video URL
        elif "youtu.be/" in url:
            video_id = url.split("youtu.be/")[1].split("?")[0]
            return extract_youtube_transcript(video_id)
        
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        return " ".join(soup.stripped_strings), None
    except Exception as e:
        return None, str(e)



def extract_youtube_transcript(video_id):
    try:
        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        # Extract the 'text' from each entry in the transcript and join them
        formatted_transcript = " ".join(entry['text'] for entry in transcript)
        return formatted_transcript, None
    except Exception as e:
        return None, str(e)
    
def transcribe_audio_with_openai(filename, openai_api_key):
    try:
        with open(filename, 'rb') as audio_file:
            response = requests.post(
                "https://api.openai.com/v1/audio/transcriptions",
                headers={
                    "Authorization": f"Bearer {openai_api_key}",
                    "Content-Type": "multipart/form-data"
                },
                data={"model": "whisper-1"},
                files={"file": audio_file}
            )
            response.raise_for_status()
            return response.json().get('text'), None
    except Exception as e:
        return None, str(e)    



if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=8080)